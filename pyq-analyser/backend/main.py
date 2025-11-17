import os
import io
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image
import google.generativeai as genai
from faiss import IndexFlatL2
import numpy as np

# -----------------------------
# Create FastAPI App
# -----------------------------
app = FastAPI()

# -----------------------------
# CORS
# -----------------------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# -----------------------------
# GOOGLE API KEY
# -----------------------------
genai.configure(api_key="enter your api key")

# -----------------------------
# LOCAL PATHS
# -----------------------------
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
POPPLER_PATH = r"C:\Users\HP\Documents\Downloads\Release-25.11.0-0\poppler-25.11.0\Library\bin"

# -----------------------------
# Vector DB
# -----------------------------
index = IndexFlatL2(768)
documents = []

# -----------------------------
# Extract text functions
# -----------------------------
def extract_text_from_pdf(file_bytes):
    images = convert_from_bytes(file_bytes, poppler_path=POPPLER_PATH)
    text = ""
    for img in images:
        text += pytesseract.image_to_string(img) + "\n"
    return text

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# -----------------------------
# Chunk text
# -----------------------------
def chunk_text(text, size=300):
    words = text.split()
    return [" ".join(words[i:i+size]) for i in range(0, len(words), size)]

# -----------------------------
# Embeddings
# -----------------------------
embed_model = "models/text-embedding-004"

def embed_text(texts):
    result = genai.embed_content(
        model=embed_model,
        content=texts,
        task_type="retrieval_document"
    )
    return np.array(result["embedding"], dtype="float32")

# -----------------------------
# ROUTE: Upload + Process Files
# -----------------------------
@app.post("/upload")
async def upload_files(file: UploadFile = File(...)):  # changed from list[UploadFile] to single UploadFile and param name to 'file'
    global documents

    file_bytes = await file.read()
    ext = file.filename.split(".")[-1].lower()

    if ext == "pdf":
        text = extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        text = extract_text_from_docx(file_bytes)
    elif ext == "pptx":
        text = extract_text_from_pptx(file_bytes)
    elif ext == "txt":
        text = file_bytes.decode("utf-8", errors="ignore")
    else:
        return {"error": f"Unsupported file format: {ext}"}

    chunks = chunk_text(text)
    embeddings = []
    for chunk in chunks:
        emb = embed_text(chunk)
        embeddings.append(emb)
        documents.append(chunk)

    index.add(np.vstack(embeddings))

    return {
        "status": "success",
        "num_chunks": len(chunks),
        "file": {"name": file.filename, "chunks": len(chunks)}
    }

# -----------------------------
# ROUTE: Ask Question
# -----------------------------
@app.post("/ask")
async def ask_question(data: dict):
    try:
        question = data.get("question", "")
        k = data.get("k", 4)

        if not question:
            return {"error": "Question is required"}

        print(f"DEBUG: Received question: {question}")

        q_emb = embed_text(question)
        print(f"DEBUG: Question embedding shape: {q_emb.shape}")

        D, I = index.search(np.array([q_emb]), k)
        retrieved = [documents[i] for i in I[0] if i < len(documents)]

        if len(retrieved) == 0:
            prompt = f"Answer the following question:\n\nQuestion: {question}"
            print("DEBUG: No context found, sending direct question prompt")
        else:
            prompt = (
                "Use the following context to answer the question. "
                "If the answer is not contained in the context, answer based on your general knowledge.\n\n"
                + "\n\n---\n\n".join(retrieved)
                + f"\n\nQuestion: {question}"
            )
            print("DEBUG: Context found, sending context prompt")

        print("DEBUG: Prompt sent to model:\n", prompt)

        model = genai.GenerativeModel("models/gemini-2.5-flash")
        response = model.generate_content(prompt)

        return {"answer": response.text}

    except Exception as e:
        print(f"ERROR in /ask: {e}")
        return {"error": "Internal server error occurred"}


